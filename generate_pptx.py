import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set slide width and height to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    DARK_GREEN = RGBColor(15, 41, 34)      # #0f2922
    GOLD = RGBColor(179, 128, 50)          # #b38032
    BG_PARCHMENT = RGBColor(247, 248, 246)  # #f7f8f6
    TEXT_DARK = RGBColor(30, 41, 59)        # #1e293b
    TEXT_MUTED = RGBColor(100, 116, 139)    # #64748b
    WHITE = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(226, 232, 240)
    ACCENT_BLUE = RGBColor(45, 127, 157)   # #2d7f9d
    ACCENT_GREEN = RGBColor(22, 163, 74)   # #16a34a

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    def add_header(slide, title_text, category_text="NILEGUARD ARCHITECTURE & TECHNICAL OVERVIEW"):
        # Top banner line
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = DARK_GREEN
        top_bar.line.color.rgb = DARK_GREEN
        
        # Category label
        tf_cat = top_bar.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = Inches(0.8)
        tf_cat.margin_top = Inches(0.15)
        
        p0 = tf_cat.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = GOLD
        p0.font.name = 'Arial'

        # Main Title
        p1 = tf_cat.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.font.name = 'Arial'

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # =========================================================
    # SLIDE 1: Title Slide
    # =========================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, DARK_GREEN)

    # Decorative shape
    deco = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    deco.fill.solid()
    deco.fill.fore_color.rgb = GOLD
    deco.line.color.rgb = GOLD

    txBox = s1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.0), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "NILEGUARD PLATFORM"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = 'Arial'

    p = tf.add_paragraph()
    p.text = "Technical Architecture & System Overview"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Arial'
    p.space_after = Pt(14)

    p = tf.add_paragraph()
    p.text = "المنظومة الوطنية الذكية للإنذار المبكر بمخاطر الجفاف الزراعي في صعيد مصر"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.font.name = 'Arial'
    p.space_after = Pt(24)

    p = tf.add_paragraph()
    p.text = "Full-Stack Architecture · FastAPI Backend · Leaflet 4D GIS · PyTorch ConvLSTM · Gemini LLM Engine"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = 'Arial'

    # Footer note
    tx_ft = s1.shapes.add_textbox(Inches(1.2), Inches(6.3), Inches(11.0), Inches(0.8))
    p_ft = tx_ft.text_frame.paragraphs[0]
    p_ft.text = "NileGuard Platform © 2026 · Arab Republic of Egypt · Ministry of Water Resources & Irrigation Alignment"
    p_ft.font.size = Pt(11)
    p_ft.font.color.rgb = TEXT_MUTED

    # =========================================================
    # SLIDE 2: System Architecture Overview
    # =========================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, BG_PARCHMENT)
    add_header(s2, "System Architecture Overview (معمارية المنظومة والربط المتكامل)")

    # Box 1: Frontend
    b1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(3.6), Inches(5.2))
    b1.fill.solid()
    b1.fill.fore_color.rgb = WHITE
    b1.line.color.rgb = ACCENT_BLUE
    b1.line.width = Pt(2)
    tf1 = b1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = Inches(0.3)
    p = tf1.paragraphs[0]
    p.text = "1. Frontend Layer (الفرونت إند)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(10)

    bullets1 = [
        "Single-Page Application (SPA) built with clean HTML5 & CSS3.",
        "Leaflet.js GIS Engine for 4D satellite basemap rendering.",
        "Chart.js time-series trajectory plotting (1958-2026).",
        "Dynamic Bilingual Engine (RTL/LTR AR/EN toggle).",
        "Agri-ROI Calculator with animated SVG liquid tank."
    ]
    for bullet in bullets1:
        p = tf1.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # Box 2: Connection / API
    b2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.6), Inches(3.6), Inches(5.2))
    b2.fill.solid()
    b2.fill.fore_color.rgb = WHITE
    b2.line.color.rgb = GOLD
    b2.line.width = Pt(2)
    tf2 = b2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = tf2.margin_top = Inches(0.3)
    p = tf2.paragraphs[0]
    p.text = "2. Integration Bridge (الربط الشبكي)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(10)

    bullets2 = [
        "Asynchronous HTTP fetch() requests from browser.",
        "JSON payload data binding (/api/forecast).",
        "Dual-Mode Resilience: Live API + Client-side Fallback Station Dictionary.",
        "Zero-delay UI updates across Leaflet markers, ranked tables, and charts."
    ]
    for bullet in bullets2:
        p = tf2.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # Box 3: Backend & AI
    b3 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.6), Inches(3.6), Inches(5.2))
    b3.fill.solid()
    b3.fill.fore_color.rgb = WHITE
    b3.line.color.rgb = ACCENT_GREEN
    b3.line.width = Pt(2)
    tf3 = b3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = tf3.margin_right = tf3.margin_top = Inches(0.3)
    p = tf3.paragraphs[0]
    p.text = "3. Backend & AI Layer (الباك إند والنمذجة)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(10)

    bullets3 = [
        "FastAPI high-performance ASGI Python server.",
        "PyTorch ConvLSTM deep learning inference pipeline.",
        "TerraClimate historical grid processing (1958-2024).",
        "Google Gemini AI Engine for contextual agronomic advice."
    ]
    for bullet in bullets3:
        p = tf3.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # =========================================================
    # SLIDE 3: Backend Architecture & FastAPI
    # =========================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, BG_PARCHMENT)
    add_header(s3, "Backend Architecture & FastAPI Server (خادم FastAPI وموديل التنبؤ)")

    # Left Column: FastAPI Details
    box_l = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = WHITE
    box_l.line.color.rgb = BORDER_COLOR
    tfl = box_l.text_frame
    tfl.word_wrap = True
    tfl.margin_left = tfl.margin_right = tfl.margin_top = Inches(0.3)

    p = tfl.paragraphs[0]
    p.text = "FastAPI Backend Implementation (api_server.py)"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(8)

    fastapi_points = [
        ("Framework", "FastAPI (Python 3.13) + Uvicorn ASGI Server running on Port 8000."),
        ("CORS Middleware", "Configured CORSMiddleware allowing cross-origin requests from GitHub Pages and localhost."),
        ("Endpoint /api/forecast", "Accepts POST/GET parameters (governorate, horizon_months). Returns JSON payload with 8 station predictions."),
        ("Deep Learning Pipeline", "Integrates PyTorch ConvLSTM model trained on TerraClimate PDSI sequence data."),
        ("Response Speed", "Inference execution latency < 45ms per query.")
    ]
    for head, body in fastapi_points:
        p = tfl.add_paragraph()
        p.text = f"• {head}: "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = GOLD
        run = p.add_run()
        run.text = body
        run.font.bold = False
        run.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # Right Column: AI Model & Inference Logic
    box_r = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.2))
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = WHITE
    box_r.line.color.rgb = BORDER_COLOR
    tfr = box_r.text_frame
    tfr.word_wrap = True
    tfr.margin_left = tfr.margin_right = tfr.margin_top = Inches(0.3)

    p = tfr.paragraphs[0]
    p.text = "Deep Learning Model & Transfer Learning Pipeline"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(8)

    model_points = [
        ("Model Type", "Spatial-Temporal ConvLSTM Neural Network."),
        ("Input Feature Grid", "TerraClimate 1958–2024 monthly Palmer Drought Severity Index (PDSI)."),
        ("Transfer Learning", "Fine-tuned model updated through end of 2025 to project 2026 drought trajectory."),
        ("Forecast Horizons", "Predicts Month +3 (Mar 2026), Month +6 (Jun 2026), and Month +9 (Sep 2026)."),
        ("Validation Metrics", "Predictive Accuracy coefficient of determination R² = 0.760.")
    ]
    for head, body in model_points:
        p = tfr.add_paragraph()
        p.text = f"• {head}: "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = ACCENT_BLUE
        run = p.add_run()
        run.text = body
        run.font.bold = False
        run.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # =========================================================
    # SLIDE 4: Frontend Architecture & UI/UX
    # =========================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, BG_PARCHMENT)
    add_header(s4, "Frontend Architecture & UI/UX Design (تصميم واجهة المستخدم والأداء)")

    box_f = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.6), Inches(5.2))
    box_f.fill.solid()
    box_f.fill.fore_color.rgb = WHITE
    box_f.line.color.rgb = BORDER_COLOR
    tff = box_f.text_frame
    tff.word_wrap = True
    tff.margin_left = tff.margin_right = tff.margin_top = Inches(0.4)

    p = tff.paragraphs[0]
    p.text = "Editorial Academic Design System (أنظمة التصميم والأداء)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(12)

    ui_features = [
        ("Design Palette", "Light Editorial Parchment theme (#f7f8f6), Royal Emerald Green (#0f2922), Muted Gold (#b38032), and Slate Text."),
        ("Typography Hierarchy", "Cairo & Amiri for Arabic serif typography; DM Serif Display & Inter for English academic layout."),
        ("Dynamic Localization Engine", "Seamless Arabic/English language toggle (updateHeroLanguage()) updating HTML dir (RTL/LTR), titles, chip labels, tables, and charts instantly without page reload."),
        ("Responsive Tab Navigation", "3 Primary Tabs: Monitoring & Forecast Dashboard, Smart Crop Intelligence, and Agri-ROI Water Savings Calculator."),
        ("Clean Non-AI Aesthetics", "Stripped of all decorative emojis and unnecessary model metrics to maintain official governmental presentation standards.")
    ]

    for title, desc in ui_features:
        p = tff.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = GOLD
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # =========================================================
    # SLIDE 5: Frontend-Backend Integration Detail
    # =========================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, BG_PARCHMENT)
    add_header(s5, "Frontend-Backend Integration (طريقة الربط وآلية البيانات)")

    # Left: Code Flow
    box_c = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.2))
    box_c.fill.solid()
    box_c.fill.fore_color.rgb = DARK_GREEN
    box_c.line.color.rgb = DARK_GREEN
    tfc = box_c.text_frame
    tfc.word_wrap = True
    tfc.margin_left = tfc.margin_right = tfc.margin_top = Inches(0.3)

    p = tfc.paragraphs[0]
    p.text = "Async Data Flow (آلية طلب البيانات والتنسيق)"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = GOLD
    p.space_after = Pt(10)

    code_steps = [
        "1. User Selects Parameters (Governorate / Horizon Month +3, +6, +9).",
        "2. Trigger runInference() JS Function.",
        "3. HTTP fetch() request dispatched to FastAPI endpoint (/api/forecast).",
        "4. Server returns JSON object containing station PDSI forecasts.",
        "5. UI Callbacks update Satellite Map Markers, Ranked Risk Table, and Trajectory Charts simultaneously."
    ]
    for step in code_steps:
        p = tfc.add_paragraph()
        p.text = step
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

    # Right: Resilience & Fallback
    box_res = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(1.6), Inches(5.2), Inches(5.2))
    box_res.fill.solid()
    box_res.fill.fore_color.rgb = WHITE
    box_res.line.color.rgb = BORDER_COLOR
    tfres = box_res.text_frame
    tfres.word_wrap = True
    tfres.margin_left = tfres.margin_right = tfres.margin_top = Inches(0.3)

    p = tfres.paragraphs[0]
    p.text = "Resilience & Fallback Architecture (مرونة المنظومة)"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(10)

    res_points = [
        ("Offline Resilience", "Includes a pre-computed client-side dictionary for 8 Upper Egypt stations (Asyut, Minya, Sohag, Qena, Luxor, Aswan, Beni Suef, Fayoum)."),
        ("Failover Mechanism", "If FastAPI backend server is unreachable, JS automatically falls back to local dataset without breaking UI."),
        ("Cache Busting Protocol", "Versioned assets (styles.css?v=30.0 & app.js?v=42.0) ensure instant updates across GitHub Pages.")
    ]
    for h, d in res_points:
        p = tfres.add_paragraph()
        p.text = f"• {h}: "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = ACCENT_BLUE
        r = p.add_run()
        r.text = d
        r.font.bold = False
        r.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # =========================================================
    # SLIDE 6: 4D GIS Satellite Map & Time Player
    # =========================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, BG_PARCHMENT)
    add_header(s6, "Interactive 4D GIS Satellite Map (خريطة الأقمار الصناعية ومحاكاة 4D)")

    box_gis = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.6), Inches(5.2))
    box_gis.fill.solid()
    box_gis.fill.fore_color.rgb = WHITE
    box_gis.line.color.rgb = BORDER_COLOR
    tfg = box_gis.text_frame
    tfg.word_wrap = True
    tfg.margin_left = tfg.margin_right = tfg.margin_top = Inches(0.4)

    p = tfg.paragraphs[0]
    p.text = "Leaflet GIS & 4D Simulation Engine (أنظمة الرصد الفضائي)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(12)

    gis_features = [
        ("Multi-Basemap Switcher", "Esri World Imagery (Satellite), OpenStreetMap Terrain, and CartoDB Dark basemaps."),
        ("4D Time Player Toolbar", "Interactive range slider (1M to 12M horizon for 2026: Jan 2026 to Dec 2026) with Play/Pause and 1x, 2x, 4x speed controls."),
        ("Dynamic Circle Markers", "Color-coded station markers with radii dynamically scaled proportional to absolute drought magnitude (|PDSI|)."),
        ("4 Official Certified Drought Categories", "Cat 0: Normal/Mild (PDSI > -1.0, Green #16a34a) | Cat 1: Moderate (-2.0 < PDSI <= -1.0, Amber #d97706) | Cat 2: Severe (-3.0 < PDSI <= -2.0, Orange #ea580c) | Cat 3: Extreme (PDSI <= -3.0, Red #b91c1c).")
    ]
    for h, d in gis_features:
        p = tfg.add_paragraph()
        p.text = f"• {h}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = GOLD
        r = p.add_run()
        r.text = d
        r.font.bold = False
        r.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # =========================================================
    # SLIDE 7: Time Series Trajectory & Forecast Chart
    # =========================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, BG_PARCHMENT)
    add_header(s7, "Time Series Trajectory & Forecast Chart (مسار الجفاف والتنبؤ المائي)")

    box_chart = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.6), Inches(5.2))
    box_chart.fill.solid()
    box_chart.fill.fore_color.rgb = WHITE
    box_chart.line.color.rgb = BORDER_COLOR
    tfch = box_chart.text_frame
    tfch.word_wrap = True
    tfch.margin_left = tfch.margin_right = tfch.margin_top = Inches(0.4)

    p = tfch.paragraphs[0]
    p.text = "Chart.js Visualization Pipeline (تحليل مسار الجفاف التاريخي والتنبؤي)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(12)

    chart_points = [
        ("Dual Dataset Binding", "Combines historical monthly observed PDSI records (1958–2024) in Steel Blue (#2d7f9d) with 2026 forecast series in Deep Orange (#ea580c)."),
        ("Clean Professional Legend", "Features clean legend title: ■ Observed (1958–2024) | ■ Forecast (completely stripped of confidence intervals per user requirements)."),
        ("Interactive Tooltips", "Hover tooltips display exact Palmer index values, date, and category classification."),
        ("Station Synchronization", "Selecting any governorate chip (Asyut, Minya, Sohag, etc.) immediately redraws the trajectory chart for that specific location.")
    ]
    for h, d in chart_points:
        p = tfch.add_paragraph()
        p.text = f"• {h}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT_BLUE
        r = p.add_run()
        r.text = d
        r.font.bold = False
        r.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # =========================================================
    # SLIDE 8: Smart Crop Intelligence Advisor
    # =========================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, BG_PARCHMENT)
    add_header(s8, "Smart Crop Intelligence Advisor (المرشد الزراعي الذكي)")

    box_crop = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.6), Inches(5.2))
    box_crop.fill.solid()
    box_crop.fill.fore_color.rgb = WHITE
    box_crop.line.color.rgb = BORDER_COLOR
    tfcrop = box_crop.text_frame
    tfcrop.word_wrap = True
    tfcrop.margin_left = tfcrop.margin_right = tfcrop.margin_top = Inches(0.4)

    p = tfcrop.paragraphs[0]
    p.text = "Targeted Agronomic Recommendations (التوصيات الزراعية الذكية)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(12)

    crop_features = [
        ("4 Certified Crop Cards", "Sakha 95 / Misr 3 Wheat, Drip Seedling Sugarcane & Sugar Beet, Manfaluti Pomegranate, and Crop Rotation / Soil Conservation."),
        ("Dynamic Agronomic Badges", "Water savings %, Planting window, Recommended irrigation method (Drip/Sprinkler), and Soil suitability."),
        ("Location-Specific Filtering", "Adapts recommendations based on soil type (Clay-Loam in Asyut, Alluvial in Sohag, Sandy-Clay in Aswan) and active PDSI severity."),
        ("Official ARC Guidance", "Aligned with Agricultural Research Center (ARC 2024) drought-resilient variety bulletins.")
    ]
    for h, d in crop_features:
        p = tfcrop.add_paragraph()
        p.text = f"• {h}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT_GREEN
        r = p.add_run()
        r.text = d
        r.font.bold = False
        r.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # =========================================================
    # SLIDE 9: Agri-ROI Calculator & Executive Bulletin
    # =========================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, BG_PARCHMENT)
    add_header(s9, "Agri-ROI Water Savings Calculator (حاسبة الوفر المائي والجدوى الاقتصادية)")

    box_roi = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.6), Inches(5.2))
    box_roi.fill.solid()
    box_roi.fill.fore_color.rgb = WHITE
    box_roi.line.color.rgb = BORDER_COLOR
    tfroi = box_roi.text_frame
    tfroi.word_wrap = True
    tfroi.margin_left = tfroi.margin_right = tfroi.margin_top = Inches(0.4)

    p = tfroi.paragraphs[0]
    p.text = "Water Savings & Economic Quantification Engine (حاسبة الوفر المائي والمالي)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(12)

    roi_points = [
        ("Interactive Feddan Slider", "Allows farmers and water engineers to select 1 to 50 feddans with real-time recalculation."),
        ("Governorate Transition Scenarios", "Custom crop conversion scenarios (e.g., Flood Wheat -> Drip Sakha 95 in Asyut; Flood Sugarcane -> Sugar Beet in Luxor)."),
        ("Animated SVG Liquid Tank", "Visualizes % water volume saved dynamically with CSS glass sheen and animated fluid waves."),
        ("Quantified Output Cards", "1. Total Annual Water Volume Saved (m³/season) | 2. Financial Energy & Pumping Savings (EGP) | 3. Heat-Stress Resilience Index (+54%)."),
        ("Official PDF Printable Bulletin", "Generates print-ready national executive drought bulletin for irrigation directorates with official document serial numbers.")
    ]
    for h, d in roi_points:
        p = tfroi.add_paragraph()
        p.text = f"• {h}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = GOLD
        r = p.add_run()
        r.text = d
        r.font.bold = False
        r.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # =========================================================
    # SLIDE 10: Official Data Sources & References
    # =========================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10, BG_PARCHMENT)
    add_header(s10, "Official Data Sources & Scientific References (المصادر والمراجع الرسمية المعتمدة)")

    box_ref = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.6), Inches(5.2))
    box_ref.fill.solid()
    box_ref.fill.fore_color.rgb = WHITE
    box_ref.line.color.rgb = BORDER_COLOR
    tfref = box_ref.text_frame
    tfref.word_wrap = True
    tfref.margin_left = tfref.margin_right = tfref.margin_top = Inches(0.4)

    p = tfref.paragraphs[0]
    p.text = "Certified Data Sources & Institutional Benchmarks (المصادر الرسمية)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(12)

    refs = [
        ("01. TerraClimate Satellite Dataset", "Monthly climate grid (1958–2024) from University of Idaho Climatology Lab for historical PDSI baseline calibration."),
        ("02. MWRI National Water Quotas", "Ministry of Water Resources & Irrigation official crop water consumption guide for Upper Egypt governorates (2022-2024)."),
        ("03. ARC Technical Bulletins", "Agricultural Research Center technical releases for drought-tolerant crop varieties (Sakha 95 wheat, drip sugarcane, beet)."),
        ("04. WMRI Irrigation Efficiency Standards", "Water Management Research Institute standards for drip efficiency gains + energy/fertilizer pricing factor (2.2 EGP/m³)."),
        ("05. FAO Irrigation Paper 56 & HCWW", "Food and Agriculture Organization ET0 standards + Holding Company for Water & Wastewater domestic benchmark (400 m³/family/year).")
    ]
    for h, d in refs:
        p = tfref.add_paragraph()
        p.text = f"• {h}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GREEN
        r = p.add_run()
        r.text = d
        r.font.bold = False
        r.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # =========================================================
    # SLIDE 11: Generative AI & Gemini LLM Integration
    # =========================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11, BG_PARCHMENT)
    add_header(s11, "Generative AI & Gemini LLM Integration (تكامل المستشار الذكي محرك Gemini)")

    box_llm = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.6), Inches(5.2))
    box_llm.fill.solid()
    box_llm.fill.fore_color.rgb = WHITE
    box_llm.line.color.rgb = BORDER_COLOR
    tfllm = box_llm.text_frame
    tfllm.word_wrap = True
    tfllm.margin_left = tfllm.margin_right = tfllm.margin_top = Inches(0.4)

    p = tfllm.paragraphs[0]
    p.text = "Floating AI Agronomic Advisor (المستشار الزراعي التوليدي)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(12)

    llm_features = [
        ("Powered by Gemini AI Engine", "Integrates Google Gemini AI Engine for real-time agronomic query processing."),
        ("Contextual Prompt Injection", "Automatically injects current active governorate, soil classification, primary agriculture, and predicted PDSI score into query context."),
        ("Floating Widget Interface", "Accessible via floating launcher button (FAB) with minimize/expand controls and active context banner."),
        ("Quick Suggestion Chips", "Interactive chips for instant queries (Wheat varieties, Water deficit management, Pomegranate irrigation, Sugarcane alternatives)."),
        ("Bilingual Intelligence", "Delivers fully localized, highly detailed agronomic guidance in both Arabic and English.")
    ]
    for h, d in llm_features:
        p = tfllm.add_paragraph()
        p.text = f"• {h}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT_BLUE
        r = p.add_run()
        r.text = d
        r.font.bold = False
        r.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # =========================================================
    # SLIDE 12: Executive Summary & Project Impact
    # =========================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12, DARK_GREEN)

    # Header
    txBox = s12.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EXECUTIVE SUMMARY & SYSTEM IMPACT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = 'Arial'

    p = tf.add_paragraph()
    p.text = "الملخص التنفيذي وأثر المنظومة القومية"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Arial'

    # Impact Grid (4 Cards)
    sdgs = [
        ("SDG 13: Climate Action", "Early drought warning systems for climate adaptation across Upper Egypt."),
        ("SDG 6: Clean Water", "Water quota optimization & 32%+ saving in agricultural canal deliveries."),
        ("SDG 2: Zero Hunger", "Protecting crop yield & guiding farmers to heat-stress resilient varieties."),
        ("SDG 15: Life on Land", "Combating desertification & soil degradation in Nile Valley farmland.")
    ]

    lefts = [Inches(0.8), Inches(6.8), Inches(0.8), Inches(6.8)]
    tops = [Inches(2.2), Inches(2.2), Inches(4.5), Inches(4.5)]

    for i in range(4):
        card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lefts[i], tops[i], Inches(5.7), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(27, 60, 52)
        card.line.color.rgb = GOLD
        card.line.width = Pt(1.5)
        tfc = card.text_frame
        tfc.word_wrap = True
        tfc.margin_left = tfc.margin_right = tfc.margin_top = Inches(0.25)
        
        p = tfc.paragraphs[0]
        p.text = sdgs[i][0]
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = GOLD
        p.space_after = Pt(4)
        
        p = tfc.add_paragraph()
        p.text = sdgs[i][1]
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

    # Footer slogan
    tx_slogan = s12.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.6))
    p_s = tx_slogan.text_frame.paragraphs[0]
    p_s.text = "« نستبق الجفاف.. لأمان النيل ونماء الأرض » · Ahead of Drought: Securing the Nile, Sustaining the Land"
    p_s.alignment = PP_ALIGN.CENTER
    p_s.font.size = Pt(13)
    p_s.font.bold = True
    p_s.font.color.rgb = GOLD

    output_path = r"C:\Users\student\.gemini\antigravity\scratch\nileguard_app\NileGuard_Technical_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_deck()
